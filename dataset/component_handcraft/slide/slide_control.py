from pptx import Presentation
from playwright.sync_api import sync_playwright
import cv2
import numpy as np
from PIL import Image
import io
import os
import time
import json
import argparse

parser = argparse.ArgumentParser(description="Process an image and draw text on it.")

parser.add_argument("--screensize_name", type=str, required=True)
args = parser.parse_args()


class PPTWebAutomation:
    def __init__(self, ppt_path, url, screensize_name):
        self.ppt_path = ppt_path
        self.url = url
        self.screensize_name = screensize_name
        self.screensize_dict = {
            "1280*720": (1280, 720),
            "1920*1080": (1920, 1080),
            "3840*2160": (3840, 2160),
        }
        self.screensize = self.screensize_dict[screensize_name]
        self.offset_dict = {
            "1280*720": {"x": 382 / 1280, "y": 188 / 720},
            "1920*1080": {"x": 445 / 1920, "y": 225 / 1080},
            "3840*2160": {"x": 638 / 3840, "y": 334 / 2160},
        }
        self.offset = self.offset_dict[screensize_name]
        self.blankpos_dict = {
            "1280*720": {"x": 310, "y": 100},
            "1920*1080": {"x": 315, "y": 150},
            "3840*2160": {"x": 412, "y": 211},
        }
        self.blankpos = self.blankpos_dict[screensize_name]
        self.emu2pixel_dict = {
            "1280*720": (72 / 914400) * (60 / 75),
            "1920*1080": (72 / 914400) * (100 / 75),
            "3840*2160": (72 / 914400) * (220.1 / 75),
        }
        self.emu2pixel = self.emu2pixel_dict[screensize_name]
        self.rotatedist_dict = {
            "1280*720": 19 / 720,
            "1920*1080": 19 / 1080,
            "3840*2160": 19 / 2160,
        }
        self.rotatedist = self.rotatedist_dict[screensize_name]
        self.ppt = Presentation(ppt_path)

    def get_textbox_bbox(self, slide):
        bboxes = []
        for shape in slide.shapes:
            bbox = {}
            bbox["rotation"] = shape.rotation

            left = shape.left * self.emu2pixel / self.screensize[0]
            top = shape.top * self.emu2pixel / self.screensize[1]
            width = shape.width * self.emu2pixel / self.screensize[0]
            height = shape.height * self.emu2pixel / self.screensize[1]
            bbox["x1"] = left
            bbox["y1"] = top
            bbox["x2"] = left + width
            bbox["y2"] = top + height

            if shape.has_text_frame:
                print("text:", shape.text)
                bbox["text"] = shape.text

            else:
                bbox["text"] = None
            if bbox["rotation"] == 0:
                bboxes.append(bbox)

        return bboxes

    def apply_offset(self, bboxes):
        offset_bboxes = []
        for bbox in bboxes:
            offset_bbox = bbox
            offset_bbox["x1"] = bbox["x1"] + self.offset["x"]
            offset_bbox["y1"] = bbox["y1"] + self.offset["y"]
            offset_bbox["x2"] = bbox["x2"] + self.offset["x"]
            offset_bbox["y2"] = bbox["y2"] + self.offset["y"]
            offset_bboxes.append(offset_bbox)
        return offset_bboxes

    def draw_bbox(self, image, bbox):
        img = np.array(image)
        start_point = (
            int(bbox["x1"] * self.screensize[0]),
            int(bbox["y1"] * self.screensize[1]),
        )
        end_point = (
            int(bbox["x2"] * self.screensize[0]),
            int(bbox["y2"] * self.screensize[1]),
        )
        color = (0, 255, 0)
        thickness = 2
        img = cv2.rectangle(img, start_point, end_point, color, thickness)
        return Image.fromarray(img)

    def process_slides(self):
        os.makedirs("slides", exist_ok=True)
        action_list = []
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=False)
            page = browser.new_page(
                viewport={"width": self.screensize[0], "height": self.screensize[1]}
            )
            page.goto(self.url)

            time.sleep(15)
            page.mouse.click(self.blankpos["x"], self.blankpos["y"])
            print("blankpos")
            page.wait_for_timeout(200)
            page.keyboard.press("Meta+-")
            page.wait_for_timeout(100)
            page.keyboard.press("Meta+-")
            page.wait_for_timeout(100)
            print("press meta+-")
            for slide_idx, slide in enumerate(self.ppt.slides):
                print(f"Processing slide {slide_idx}")

                bboxes = self.get_textbox_bbox(slide)
                offset_bboxes = self.apply_offset(bboxes)
                print(f"Found {len(offset_bboxes)} text boxes: {str(offset_bboxes)}")

                page.keyboard.press("Meta+A")
                page.wait_for_timeout(500)
                print("meta+A")

                output_dir = f"slides_{self.screensize_name}/slide_{slide_idx}"
                os.makedirs(output_dir, exist_ok=True)

                with open(os.path.join(output_dir, "bbox.json"), "w") as file:
                    json.dump(offset_bboxes, file)
                screenshot = page.screenshot()
                print("take screenshot")
                image = Image.open(io.BytesIO(screenshot))
                image.save(f"{output_dir}/original.png")

                for bbox_idx, bbox in enumerate(offset_bboxes):

                    marked_image = self.draw_bbox(image, bbox)

                marked_image.save(f"{output_dir}/annotated.png")

                for bbox_idx, bbox in enumerate(offset_bboxes):
                    page.keyboard.press("Tab")
                    page.wait_for_timeout(200)
                    screenshot = page.screenshot()
                    print("take screenshot")
                    image = Image.open(io.BytesIO(screenshot))
                    image.save(f"{output_dir}/original_bbox_{bbox_idx}.png")
                    os.makedirs(output_dir, exist_ok=True)
                    marked_image = self.draw_bbox(image, bbox)
                    marked_image.save(f"{output_dir}/annotated_bbox_{bbox_idx}.png")

                if slide_idx < len(self.ppt.slides) - 1:
                    page.wait_for_timeout(500)
                    page.mouse.click(self.blankpos["x"], self.blankpos["y"])
                    print("clicking blank pos")
                    page.wait_for_timeout(500)
                    page.keyboard.press("ArrowDown")
                    print("press arrow down")
                    page.wait_for_timeout(1000)

            browser.close()


if __name__ == "__main__":
    automation = PPTWebAutomation(
        ppt_path="<path to your pptx file>",
        url="<url of the pptx file>",
        screensize_name=args.screensize_name,
    )
    automation.process_slides()
